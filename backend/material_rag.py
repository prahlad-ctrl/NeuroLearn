"""
RAG pipeline for uploaded PDF/PPTX materials.

- Extracts text from uploaded files (PDF, PPTX, TXT)
- Cleans and normalises extracted text
- Chunks with sentence-aware splitting + sliding-window overlap
- Stores TF-IDF vectors in a lightweight in-memory store (per-session)
- Retrieves relevant chunks via cosine similarity with score thresholds
- Supports multi-file uploads per session (additive chunk store)

Uses scikit-learn TfidfVectorizer for lightweight vector search (no GPU needed).
"""

from __future__ import annotations

import io
import re
import math
from typing import BinaryIO

from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def extract_text_pdf(file: BinaryIO) -> str:
    """Extract text from a PDF file using PyPDF2."""
    import PyPDF2
    reader = PyPDF2.PdfReader(file)
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n\n".join(pages)


def extract_text_pptx(file: BinaryIO) -> str:
    """Extract text from a PPTX file, preserving slide structure."""
    from pptx import Presentation
    prs = Presentation(file)
    slides: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
            # Also extract text from tables
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        if parts:
            slides.append(f"[Slide {slide_num}]\n" + "\n".join(parts))
    return "\n\n".join(slides)


def extract_text(filename: str, file: BinaryIO) -> str:
    """Route to the correct extractor based on file extension."""
    lower = filename.lower()
    if lower.endswith(".pdf"):
        raw = extract_text_pdf(file)
    elif lower.endswith(".pptx"):
        raw = extract_text_pptx(file)
    else:
        data = file.read()
        raw = data.decode("utf-8", errors="ignore") if isinstance(data, bytes) else data

    return _clean_text(raw)


def _clean_text(text: str) -> str:
    """Normalise whitespace, remove control chars, collapse blank lines."""
    # Remove non-printable chars except newlines/tabs
    text = re.sub(r"[^\S\n\t]+", " ", text)
    # Collapse 3+ newlines into 2
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Strip each line
    lines = [line.strip() for line in text.split("\n")]
    return "\n".join(lines).strip()


# ---------------------------------------------------------------------------
# Sentence-aware chunking
# ---------------------------------------------------------------------------

_SENTENCE_SPLIT = re.compile(r"(?<=[.!?])\s+|\n{2,}")


def chunk_text(text: str, max_tokens: int = 500, overlap: int = 80) -> list[str]:
    """
    Split text into chunks of roughly `max_tokens` words.

    Strategy:
    1. Split into sentences / paragraphs first
    2. Greedily pack sentences into chunks up to max_tokens
    3. Overlap by recycling the last `overlap` words into the next chunk

    This avoids cutting mid-sentence and produces more coherent RAG context.
    """
    if not text.strip():
        return []

    # Split into sentence-like segments
    segments = _SENTENCE_SPLIT.split(text)
    segments = [s.strip() for s in segments if s.strip()]

    # If very short, return as single chunk
    total_words = sum(len(s.split()) for s in segments)
    if total_words <= max_tokens:
        return [text.strip()]

    chunks: list[str] = []
    current_words: list[str] = []
    current_len = 0

    for segment in segments:
        seg_words = segment.split()
        seg_len = len(seg_words)

        # If adding this segment would exceed limit, flush
        if current_len + seg_len > max_tokens and current_words:
            chunks.append(" ".join(current_words))
            # Keep last `overlap` words for context continuity
            if overlap > 0 and len(current_words) > overlap:
                current_words = current_words[-overlap:]
                current_len = len(current_words)
            else:
                current_words = []
                current_len = 0

        current_words.extend(seg_words)
        current_len += seg_len

    # Flush remaining
    if current_words:
        chunk = " ".join(current_words)
        # Avoid tiny trailing chunks — merge with previous if too small
        if chunks and current_len < max_tokens // 4:
            chunks[-1] += " " + chunk
        else:
            chunks.append(chunk)

    return chunks


# ---------------------------------------------------------------------------
# In-memory vector store (per-session, supports multi-file uploads)
# ---------------------------------------------------------------------------

# session_id -> {"chunks": [...], "vectorizer": TfidfVectorizer, "matrix": sparse, "filenames": [...]}
_stores: dict[str, dict] = {}


def store_chunks(session_id: str, chunks: list[str], filename: str = "") -> int:
    """
    Vectorize and store chunks for a session. Additive — uploading a second
    file appends chunks rather than replacing.
    Returns total chunk count for this session.
    """
    if not chunks:
        return len(_stores.get(session_id, {}).get("chunks", []))

    existing = _stores.get(session_id)
    if existing:
        all_chunks = existing["chunks"] + chunks
        filenames = existing.get("filenames", []) + [filename] * len(chunks)
    else:
        all_chunks = chunks
        filenames = [filename] * len(chunks)

    vectorizer = TfidfVectorizer(
        stop_words="english",
        max_features=8000,
        ngram_range=(1, 2),        # unigrams + bigrams for better matching
        sublinear_tf=True,         # dampen high-frequency terms
    )
    matrix = vectorizer.fit_transform(all_chunks)

    _stores[session_id] = {
        "chunks": all_chunks,
        "vectorizer": vectorizer,
        "matrix": matrix,
        "filenames": filenames,
    }
    return len(all_chunks)


def retrieve_chunks(
    session_id: str,
    query: str,
    top_k: int = 5,
    min_score: float = 0.05,
) -> list[str]:
    """
    Return the top-k most relevant chunks for a query.
    Filters out chunks below `min_score` cosine similarity.
    """
    store = _stores.get(session_id)
    if not store:
        return []

    query_vec = store["vectorizer"].transform([query])
    scores = cosine_similarity(query_vec, store["matrix"]).flatten()

    # Sort by score descending, take top_k that pass threshold
    ranked = scores.argsort()[::-1]
    results: list[str] = []
    for idx in ranked:
        if len(results) >= top_k:
            break
        if scores[idx] >= min_score:
            results.append(store["chunks"][idx])

    # Fallback: if nothing passed the threshold, return top 2 anyway
    if not results:
        for idx in ranked[:2]:
            results.append(store["chunks"][idx])

    return results


def get_material_stats(session_id: str) -> dict | None:
    """Return stats about stored material for a session."""
    store = _stores.get(session_id)
    if not store:
        return None
    unique_files = set(f for f in store.get("filenames", []) if f)
    total_words = sum(len(c.split()) for c in store["chunks"])
    return {
        "chunks": len(store["chunks"]),
        "files": list(unique_files),
        "total_words": total_words,
    }


def has_material(session_id: str) -> bool:
    return session_id in _stores


def clear_material(session_id: str):
    _stores.pop(session_id, None)


# ---------------------------------------------------------------------------
# Prompt builders for RAG-based generation
# ---------------------------------------------------------------------------

def build_rag_lesson_prompt(chunks: list[str], subject: str, level: str) -> str:
    context = "\n\n---\n\n".join(chunks)
    return (
        f"You are a {subject} tutor. Using ONLY the study material below, "
        f"generate a structured lesson suitable for a {level}-level student.\n\n"
        "Requirements:\n"
        "- Use clear headings (## style) to organise topics\n"
        "- Explain concepts in the student's own level of understanding\n"
        "- Include at least one worked example from the material\n"
        "- Highlight key definitions and formulas\n"
        "- End with 3 key takeaways\n\n"
        f"STUDY MATERIAL:\n{context}\n\n"
        "Generate the lesson now."
    )


def build_rag_exercise_prompt(
    chunks: list[str], subject: str, level: str, question_type: str = "mixed"
) -> str:
    context = "\n\n---\n\n".join(chunks)
    type_instruction = _type_format_instruction(question_type)
    return (
        f"You are a {subject} tutor. Using ONLY the study material below, "
        f"generate 5 practice questions for a {level}-level student.\n\n"
        "Requirements:\n"
        "- Questions must be answerable from the provided material\n"
        "- Cover different sections of the material\n"
        "- Progress from easier to harder\n\n"
        f"Question format:\n{type_instruction}\n\n"
        "Return ONLY a JSON array with no extra text.\n\n"
        f"STUDY MATERIAL:\n{context}\n\n"
        "Generate the questions now."
    )


def _type_format_instruction(question_type: str) -> str:
    if question_type == "mcq":
        return (
            'Each question must be MCQ format with keys: "type", "question", '
            '"options" (array of 4 strings), "answer" (the correct option text).'
        )
    elif question_type == "true_false":
        return (
            'Each question must be True/False format with keys: "type", "question", '
            '"answer" (boolean true or false).'
        )
    elif question_type == "short":
        return (
            'Each question must be short-answer format with keys: "type", "question", '
            '"answer" (brief text).'
        )
    elif question_type == "qa":
        return (
            'Each question must be descriptive Q&A format with keys: "type", "question", '
            '"expected_points" (array of key points the answer should cover).'
        )
    else:
        return (
            "Use a mix of question types. Each question must have a \"type\" field "
            "(one of: \"mcq\", \"true_false\", \"short\", \"qa\") and appropriate fields:\n"
            "- mcq: question, options (array of 4), answer\n"
            "- true_false: question, answer (boolean)\n"
            "- short: question, answer (text)\n"
            "- qa: question, expected_points (array of strings)"
        )
